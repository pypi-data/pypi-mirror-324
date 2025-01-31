from pptx.util import Pt        
import re 
from pptx.dml.color import RGBColor
from SlideSlide.Config import Colors
from datetime import datetime

def get_current_time():
    return datetime.now().strftime("%Y_%m_%d_%H_%M_%S")

def add_textbox(slide , text , font , font_size , font_color ,left, top ,width , height):
    title_box = slide.shapes.add_textbox(Pt(left),Pt(top),Pt(width),Pt(height))
    text_frame = title_box.text_frame
    text_frame.word_wrap=True
    
    p = text_frame.add_paragraph()
    vals = extract_bold_noraml_text_ordered(text)
    for idx , data in enumerate(vals):
        type , text = data
        run = p.add_run()
        run.text=text 
        try:
            run.font.name=font 
        except:...
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        if type == 'bold':
            run.font.bold=True

def add_content_textbox(
        slide ,
        text , 
        font,
        font_size , 
        font_color,
        left , 
        top , 
        width , 
        height 
):
    content_box = slide.shapes.add_textbox(Pt(left),Pt(top),Pt(width),Pt(height))    
    text_frame = content_box.text_frame
    text_frame.word_wrap=True

    p = text_frame.add_paragraph()
    vals = extract_bold_noraml_text_ordered(text)
    for idx , data in enumerate(vals):
        type , text = data
        run = p.add_run()
        run.text=text 
        try:
            run.font.name=font 
        except:...
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        if type == 'bold':
            run.font.bold=True
     

def add_brandname_textbox(
        slide ,
        text , 
        font,
        bold,
        font_size , 
        font_color,
        left , 
        top , 
        width , 
        height 
):
    content_box = slide.shapes.add_textbox(Pt(left),Pt(top),Pt(width),Pt(height))    
    text_frame = content_box.text_frame
    text_frame.word_wrap=True

    p = text_frame.add_paragraph()
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(font_size)  
    try:
        run.font.name=font 
    except:...
    run.font.color.rgb = font_color
    run.font.bold = bold
     


def extract_bold_noraml_text_ordered(input_string:str):
    pattern = r'(\*\*(.*?)\*\*|[^*]+)'
    parts = re.findall(pattern,input_string)
    result=[]
    for part in parts:
        if part[0].startswith('**'):
            result.append(('bold',part[0][2:-2]))
        else:
            result.append(('normal',part[0]))
    return result