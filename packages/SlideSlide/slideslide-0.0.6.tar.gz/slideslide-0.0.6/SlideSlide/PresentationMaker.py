from pptx import Presentation 
from pptx.util import Inches , Pt
from pptx.dml.color import RGBColor
from SlideSlide.utils import add_textbox , add_content_textbox , add_brandname_textbox , get_current_time
from SlideSlide.Config import Colors
import os 
import json


def MakePresentation(
        presentation_content:list[dict],
        presentation_name:str=None,
        font:str=None,
        add_ending_slide:bool=True,
        template_name:str='SapphireBlue',
        custom_template:str=None,
        template_mode:str='light',
        brand_name:str=None,
        verbose:bool=False
        ):
    
    """
    Make presentation from the input json 

    Parameters:
    presentation_content : list[dict] -> {"title":"Title of Slide" , "content":"Content for slide"}
    presentation_name : str -> Name of presentation
    add_ending_slide : bool -> Adds ending slide with title thaks for your time
    template_name : str -> name of template
    ustom_template : str -> custom json template
    template_mode : str -> light / dark
    brand_name : str -> brand name 
    verbose : bool -> shows log message on console

    Returns:
    Saves presentaion at root folder
    """

    file_path = os.path.join(os.path.dirname(__file__))

    if custom_template:
        if verbose:print(f"[INFO] Loading Custom Template")
        is_custom_template_exist=os.path.exists(custom_template)
        if not is_custom_template_exist:
            if verbose:print(f"[INFO] Custom Template File Path error using default tempalte 'SapphireBlue'")
            with open(f"{file_path}/Themes/SapphireBlue","r") as f:
                theme=json.load(f) 
        else:
            file_path=""
            with open(f"{custom_template}","r") as f:
                theme=json.load(f)
    else:
        if verbose:print("[INFO] Loading Themes") 
        themes=os.listdir(f"{file_path}/Themes")
        if f"{template_name}.json" in themes:
            with open(f"{file_path}/Themes/{template_name}.json",'r') as f:
                theme=json.load(f)
            if verbose:print(f"[INFO] Using {template_name} Theme")
        else:
            with open(f"{file_path}/Themes/SapphireBlue.json",'r') as f:
                theme=json.load(f)
            if verbose:print("[INFO] Theme Not Found Using Default SapphireBlue Theme")
   
    presentaion=Presentation()
    presentaion.slide_width=Inches(1920//96)
    presentaion.slide_height=Inches(1080//96)

    if verbose:print("[INFO] Presentation Size 1920x1080")

    for slide_number , slide_content in enumerate(presentation_content):
        # addign a blank slide no 6 is blank one
        slide = presentaion.slides.add_slide(presentaion.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid() 
        fill.fore_color.rgb = Colors.SAFEWHITE.value if template_mode=='light' else Colors.SAFEBLACK.value

        # +------------------------------------------------------------------------------------------------+
        # +---While Updating code do not add if else in fuction args this is doable but makes code messy---+
        # +------------------------------------------------------------------------------------------------+

        if slide_number==0:
            # this is cover slide
            coverslide = theme['CoverSlide']
            
            image_path=f"{file_path}{coverslide['Image']}"
            if os.path.exists(image_path):
                slide.shapes.add_picture(
                    image_path,
                    left=Inches(coverslide["ImageLeft"]),
                    top=Inches(coverslide["ImageTop"]),
                    width=Inches(coverslide["ImageWidth"]),
                    height=Inches(coverslide["ImageHeight"]),
                )
            else:
                print(f"[INOF] Image path for coverslide does not exist")

            add_textbox(
                slide=slide , 
                text=slide_content['title'],
                font=font,
                font_size=coverslide['TitleFontSize'],
                font_color=Colors.SAFEWHITE.value if slide_number==0 else Colors.SAFEBLACK.value,
                left=coverslide['TitleLeft'],
                top=coverslide["TitleTop"],
                width=coverslide["TitleWidth"],
                height=coverslide["TitleHeight"]
                )
            add_content_textbox(
                slide=slide ,
                text = slide_content['content'].strip(),
                font=font,
                font_size=coverslide['ContentFontSize'],
                font_color=Colors.SAFEWHITE.value if slide_number==0 else Colors.SAFEBLACK.value,
                top=coverslide['ContentTop'],
                left=coverslide['ContentLeft'],
                width=coverslide['ContentWidth'],
                height=coverslide['ContentHeight']
            )
            if brand_name is not None:
                add_brandname_textbox(
                    slide=slide,
                    text = brand_name,
                    font=font,
                    bold=True,
                    font_size = 22,
                    font_color=Colors.SAFEWHITE.value,
                    top=700,
                    left=1000,
                    width=400,
                    height=26
                )
        else:
            #This is ContentSlide 
            content_slide = theme['ContentSlide']
            image_path=f"{file_path}{content_slide['Image']}"
            if os.path.exists(image_path):
                slide.shapes.add_picture(
                    image_path,
                    left=Inches(content_slide["ImageLeft"]),
                    top=Inches(content_slide["ImageTop"]),
                    width=Inches(content_slide["ImageWidth"]),
                    height=Inches(content_slide["ImageHeight"]),
                )
            else:
                print("[INOF] Image path for Content slide does not exist")
            if template_mode=='light':
                font_color = Colors.SAFEBLACK.value
            else:
                font_color = Colors.SAFEWHITE.value
            add_textbox(
                slide=slide , 
                text=slide_content['title'],
                font=font,
                font_size=content_slide['TitleFontSize'],
                font_color=font_color,
                left=coverslide['TitleLeft'],
                top=coverslide["TitleTop"],
                width=coverslide["TitleWidth"],
                height=coverslide["TitleHeight"]
                )
            add_content_textbox(
                slide=slide ,
                text = slide_content['content'].strip(),
                font=font,
                font_size=content_slide['ContentFontSize'],
                font_color=font_color,
                top=content_slide['ContentTop'],
                left=content_slide['ContentLeft'],
                width=content_slide['ContentWidth'],
                height=content_slide['ContentHeight']
            )
            if brand_name is not None:
                add_brandname_textbox(
                    slide=slide,
                    text = brand_name,
                    font=font,
                    bold=True,
                    font_size = 22,
                    font_color=font_color,
                    top=670,
                    left=60,
                    width=400,
                    height=26
                )
    if add_ending_slide:
            slide = presentaion.slides.add_slide(presentaion.slide_layouts[6])
            background = slide.background
            fill = background.fill
            fill.solid() 
            fill.fore_color.rgb = Colors.SAFEWHITE.value if template_mode=='light' else Colors.SAFEBLACK.value
            coverslide = theme['CoverSlide']
            image_path=f"{file_path}{coverslide['Image']}"
            if os.path.exists(image_path):
                slide.shapes.add_picture(
                    image_path,
                    left=Inches(coverslide["ImageLeft"]),
                    top=Inches(coverslide["ImageTop"]),
                    width=Inches(coverslide["ImageWidth"]),
                    height=Inches(coverslide["ImageHeight"]),
                )
            else:
                print(f"[INOF] Image path for ending slide does not exist")
            add_textbox(
                slide=slide , 
                text="Thanks for your Time",
                font=font,
                font_size=coverslide['TitleFontSize'],
                font_color=Colors.SAFEWHITE.value ,
                left=150,
                top=200,
                width=1000,
                height=38
                )
    if presentation_name is None:
        presentation_name = get_current_time()
    presentaion.save(f"{presentation_name}.pptx")
    if verbose:print(f"[INFO] Presentation saved as {presentation_name}.pptx")
